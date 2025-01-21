import csv
import os
import openpyxl
from zhipuai import ZhipuAI
import json
import os
from langchain.text_splitter import RecursiveCharacterTextSplitter


class PDFFileLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def getParagraphs(self):
        # 示例：使用 PyPDF2 或 pdfminer 加载 PDF
        from PyPDF2 import PdfReader
        reader = PdfReader(self.file_path)
        paragraphs = [page.extract_text() for page in reader.pages]
        return paragraphs

class DocFileLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def getParagraphs(self):
        # 示例：使用 python-docx 加载 DOC 文件
        import docx
        doc = docx.Document(self.file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return paragraphs

class PPTFileLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def getParagraphs(self):
        # 示例：使用 python-pptx 加载 PPT 文件
        from pptx import Presentation
        presentation = Presentation(self.file_path)
        paragraphs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    paragraphs.extend([p.text for p in shape.text_frame.paragraphs if p.text.strip()])
        return paragraphs

class CSVFileLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def getParagraphs(self):
        with open(self.file_path, 'r', encoding='utf-8') as file:
            reader = csv.DictReader(file)
            paragraphs = []
            for row in reader:
                paragraph = ""
                for column_name, value in row.items():
                    paragraph += f"{column_name}: {value}\n"
                paragraphs.append(paragraph.strip())
            return paragraphs

class XLSXFileLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def getParagraphs(self):
        wb = openpyxl.load_workbook(self.file_path)
        ws = wb.active
        paragraphs = []
        header = [cell.value for cell in ws[1]]
        for row in ws.iter_rows(min_row=2, values_only=True):
            paragraph = ""
            for column_name, value in zip(header, row):
                paragraph += f"{column_name}: {value}\n"
            paragraphs.append(paragraph.strip())
        return paragraphs


ZHIPUAI_API_KEY = "1288e6b562440ad7a5c9bc511746dcf0.8fXJ2UA3IBm9TgXO" 

client = ZhipuAI(
    api_key=ZHIPUAI_API_KEY
)

def get_embeddings(texts, model="embedding-3"):
    """
    调用 ZhipuAI 嵌入模型生成文本嵌入向量。
    :param texts: 待生成嵌入的文本列表
    :param model: 嵌入模型名称，默认为 'embedding-3'
    :return: 嵌入向量列表
    """
    try:
        # 假设 client.embeddings.create 是正确的函数调用方式
        response = client.embeddings.create(model=model, input=texts)
        
        # 检查 response 是否有 data 属性
        if hasattr(response, 'data'):
            # 提取嵌入向量
            embeddings = [item.embedding for item in response.data]
            return embeddings
        else:
            print(f"Unexpected response structure: {response}")
            return None
    except Exception as e:
        print(f"Error generating embeddings: {e}")
        return None

def extract_main_info(file_text, model="glm-4-air", temperature=0.7):
    """调用大模型提取主要信息"""
    prompt = f"从以下内容中提取最重要的信息，过滤掉无关信息并尽量全面：\n\n{file_text}\n\n输出："
    response = client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        temperature=temperature
    )
    return response.choices[0].message.content

def generate_summary(extracted_info, model="glm-4-air", temperature=0.7):
    """调用大模型生成摘要"""
    prompt = f"对以下内容生成简洁的摘要，保留关键点：\n\n{extracted_info}\n\n输出："
    response = client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        temperature=temperature
    )
    return response.choices[0].message.content

def chunk_text(text, chunk_size=1000, overlap=200):
    """
    将长文本切分成较小的块。
    :param text: 原始文本
    :param chunk_size: 每个块的最大长度
    :param overlap: 相邻块的重叠部分长度
    :return: 切分后的文本块列表
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, chunk_overlap=overlap, separators=["\n", ".", " "]
    )
    return splitter.split_text(text)

def process_chunk_with_ai(chunk_text):
    """对单个文本切片提取信息、生成摘要并存储到数据库"""
    # 提取主要信息
    extracted_info = extract_main_info(chunk_text)

    # 生成摘要
    summary = generate_summary(extracted_info)

    # 返回提取信息和摘要
    return {"summary": summary, "extracted_info": extracted_info}



def process_files(directory):
    """
    遍历给定目录下的所有文件，提取段落信息并将结果保存到一个文本文件中。

    :param directory: 需要遍历的目录路径
    :param output_file: 保存结果的输出文件路径
    """
    file_info = []

    for root, _, files in os.walk(directory):
        for file_name in files:
            file_path = os.path.join(root, file_name)
            file_ext = os.path.splitext(file_name)[1].lower()
    
            loaders = {
                ".pdf": PDFFileLoader,
                ".doc": DocFileLoader,
                ".docx": DocFileLoader,
                ".ppt": PPTFileLoader,
                ".pptx": PPTFileLoader,
                '.csv': CSVFileLoader,
                '.xlsx': XLSXFileLoader,
            }

            if file_ext in loaders:
                try:
                    # 加载文件内容
                    paragraphs = loaders[file_ext](file_path).getParagraphs()
                    file_text = "\n".join(paragraphs)

                    # 对文件内容进行切片
                    chunks = chunk_text(file_text, chunk_size=1000, overlap=200)
                    print(f"File {file_name} split into {len(chunks)} chunks")

                    for i, chunk in enumerate(chunks):
                        try:
                            result = process_chunk_with_ai(chunk)
                            print(result)
                            
                            # 构建输出文件路径
                            output_dir = os.path.join('/home/student/zlh/GuardBot-main/Graph_RAG/ragtest/input')
                            if not os.path.exists(output_dir):
                                os.makedirs(output_dir)  # 确保目录存在
                            output_file = os.path.join(output_dir , root.split('/')[-1] + os.path.splitext(file_name)[0].lower()+str(i + 1) + '.txt')
                            print(output_file)
                            # 将所有文件信息保存到一个文本文件中
                            with open(output_file, 'w', encoding='utf-8') as f:
                                f.write(result['summary'])
                                f.write('\n')
                                f.write(result["extracted_info"])
                                f.write("\n")
                            
                            print(f"Successfully processed chunk {i + 1}/{len(chunks)} for {file_name}")
                        except Exception as e:
                            print(f"Error processing chunk {i + 1}/{len(chunks)} for {file_name}: {e}")

                except Exception as e:
                    print(f"Error processing {file_name}: {e}")
            else:
                print(f"Unsupported file type: {file_name}")




# 使用示例
if __name__ == "__main__":
    directory_to_process = "/home/student/zlh/GuardBot-main/teacher_db"  # 替换为你的目录路径
    dirs = []
    for file in os.scandir(directory_to_process):
        if file.is_dir():
            print(file)
            process_files(file)
    process_files('/home/student/zlh/GuardBot-main/teacher_db')
