import csv
import re
import openpyxl
from pypinyin import Style, lazy_pinyin
from zhipuai import ZhipuAI
import json
import os
from chromadb.config import Settings
from chromadb import PersistentClient
from sentence_transformers import SentenceTransformer
from langchain.text_splitter import RecursiveCharacterTextSplitter

ZHIPUAI_API_KEY = "1288e6b562440ad7a5c9bc511746dcf0.8fXJ2UA3IBm9TgXO" 

client = ZhipuAI(
    api_key=ZHIPUAI_API_KEY
)

embedding_model = SentenceTransformer("/home/student/zlh/model/BAAI/bge-large-zh-v1___5")

def get_embeddings(texts, model="embedding-3"):
    """
    调用 ZhipuAI 嵌入模型生成文本嵌入向量。
    :param texts: 待生成嵌入的文本列表
    :param model: 嵌入模型名称，默认为 'embedding-3'
    :return: 嵌入向量列表
    """
    try:
        # 假设 client.embeddings.create 是正确的函数调用方式
        response = client.embeddings.create(model=model, input=texts)
        
        # 检查 response 是否有 data 属性
        if hasattr(response, 'data'):
            # 提取嵌入向量
            embeddings = [item.embedding for item in response.data]
            return embeddings
        else:
            print(f"Unexpected response structure: {response}")
            return None
    except Exception as e:
        print(f"Error generating embeddings: {e}")
        return None

def chinese_to_valid_name(chinese_name):
    """将中文转为合法的集合名称"""
    # 中文转拼音
    pinyin_name = ''.join(lazy_pinyin(chinese_name, style=Style.NORMAL))
    # 替换非法字符为下划线
    valid_name = re.sub(r"[^a-zA-Z0-9_\-]", "_", pinyin_name)
    # 确保长度符合要求（3-63个字符）
    valid_name = valid_name[:63].strip("_")
    # 如果长度不足，使用默认名称
    if len(valid_name) < 3:
        valid_name = "default_collection"
    return valid_name

def chunk_text(text, chunk_size=1000, overlap=200):
    """
    将长文本切分成较小的块。
    :param text: 原始文本
    :param chunk_size: 每个块的最大长度
    :param overlap: 相邻块的重叠部分长度
    :return: 切分后的文本块列表
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, chunk_overlap=overlap, separators=["\n", ".", " "]
    )
    return splitter.split_text(text)

def extract_main_info(file_text, model="glm-4-air", temperature=0.7):
    """调用大模型提取主要信息"""
    prompt = f"从以下内容中提取最重要的信息，过滤掉无关信息并尽量全面：\n\n{file_text}\n\n输出："
    response = client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        temperature=temperature
    )
    return response.choices[0].message.content

def generate_summary(extracted_info, model="glm-4-air", temperature=0.7):
    """调用大模型生成摘要"""
    prompt = f"对以下内容生成简洁的摘要，保留关键点：\n\n{extracted_info}\n\n输出："
    response = client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        temperature=temperature
    )
    return response.choices[0].message.content

def process_file_with_ai(loader, file_path, embedding_fn, vector_db):
    """对单个文件提取信息、生成摘要并存储到数据库"""
    # 提取段落
    paragraphs = loader(file_path).getParagraphs()

    # 合并文本段落
    file_text = "\n".join(paragraphs)

    # 调用大模型提取主要信息
    extracted_info = extract_main_info(file_text)

    # 调用大模型生成摘要
    summary = generate_summary(extracted_info)

    # 将摘要存入向量数据库
    vector_db.add_documents([summary])

    # 返回提取信息和摘要
    return {"summary": summary, "extracted_info": extracted_info}

def process_chunk_with_ai(chunk_text, chunk_id, embedding_fn, vector_db):
    """对单个文本切片提取信息、生成摘要并存储到数据库"""
    # 提取主要信息
    extracted_info = extract_main_info(chunk_text)

    # 生成摘要
    summary = generate_summary(extracted_info)

    # 将摘要和提取信息存入向量数据库
    vector_db.add_documents([{"id": chunk_id, "summary": summary, "extracted_info": extracted_info}])

    # 返回提取信息和摘要
    return {"summary": summary, "extracted_info": extracted_info}

class MyVectorDBConnector:
    def __init__(self, default_collection_name, embedding_fn, db_path="/home/student/zlh/GuardBot-main/chroma_data"):
        try:
            # 转换集合名称为合法名称
            valid_collection_name = self.chinese_to_valid_name(default_collection_name)
            # 初始化客户端
            if not os.path.exists(db_path):
                os.makedirs(db_path)
            self.chroma_client = PersistentClient(path=db_path)
            self.embedding_fn = embedding_fn
            # 默认集合
            self.default_collection = self.chroma_client.get_or_create_collection(name=valid_collection_name)
            self.current_collection = self.default_collection
        except Exception as e:
            print(f"An error occurred during initialization: {e}")
            self.current_collection = None

    def chinese_to_valid_name(self,chinese_name):
        """将中文转为合法的集合名称"""
        # 中文转拼音
        pinyin_name = ''.join(lazy_pinyin(chinese_name, style=Style.NORMAL))
        # 替换非法字符为下划线
        valid_name = re.sub(r"[^a-zA-Z0-9_\-]", "_", pinyin_name)
        # 确保长度符合要求（3-63个字符）
        valid_name = valid_name[:63].strip("_")
        # 如果长度不足，使用默认名称
        if len(valid_name) < 3:
            valid_name = "default_collection"
        return valid_name

    def switch_collection(self, collection_name):
        '''切换到指定 collection'''
        try:
            self.current_collection = self.chroma_client.get_or_create_collection(name=collection_name)
            print(f"Switched to collection: {collection_name}")
        except Exception as e:
            print(f"Error switching to collection {collection_name}: {e}")

    def add_documents(self, documents):
        '''向当前 collection 添加文档与向量'''
        try:
            embeddings = self.embedding_fn([doc["summary"] for doc in documents])  # 基于 summary 生成向量
            self.current_collection.add(
                embeddings=embeddings,  # 每个文档的向量
                documents=[doc["summary"] for doc in documents],  # 存储的文档是 summary
                ids=[doc["id"] for doc in documents],  # 每个文档的 id
                metadatas=[{"extracted_info": doc["extracted_info"]} for doc in documents]  # 元数据存储 extracted_info
            )
            print("Documents added successfully.")
        except Exception as e:
            print(f"Error adding documents: {e}")


    def search(self, query, top_n=5):
        """查询向量数据库并返回完整信息"""
        try:
            results = self.current_collection.query(
                query_embeddings=self.embedding_fn([query]),
                query_texts=[query],
                n_results=top_n
            )
            detailed_results = [
                {
                    "summary": doc,
                    "extracted_info": metadata["extracted_info"]
                }
                for doc, metadata in zip(results["documents"], results["metadatas"])
            ]
            return detailed_results
        except Exception as e:
            print(f"Error during search: {e}")
            return None


    def q_search(self, query, top_n):
        '''在当前 collection 中快速检索'''
        try:
            results = self.current_collection.query(
                query_embeddings=self.embedding_fn([query]),
                query_texts=[query],
                n_results=top_n
            )
            detailed_results = [
                {
                    "summary": doc,
                    "extracted_info": metadata["extracted_info"]
                }
                for doc, metadata in zip(results["documents"], results["metadatas"])
            ]
            return detailed_results
        except Exception as e:
            print(f"Error during quick search: {e}")
            return None

class PDFFileLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def getParagraphs(self):
        # 示例：使用 PyPDF2 或 pdfminer 加载 PDF
        from PyPDF2 import PdfReader
        reader = PdfReader(self.file_path)
        paragraphs = [page.extract_text() for page in reader.pages]
        return paragraphs

class DocFileLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def getParagraphs(self):
        # 示例：使用 python-docx 加载 DOC 文件
        import docx
        doc = docx.Document(self.file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return paragraphs

class PPTFileLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def getParagraphs(self):
        # 示例：使用 python-pptx 加载 PPT 文件
        from pptx import Presentation
        presentation = Presentation(self.file_path)
        paragraphs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    paragraphs.extend([p.text for p in shape.text_frame.paragraphs if p.text.strip()])
        return paragraphs

class CSVFileLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def getParagraphs(self):
        with open(self.file_path, 'r', encoding='utf-8') as file:
            reader = csv.DictReader(file)
            paragraphs = []
            for row in reader:
                paragraph = ""
                for column_name, value in row.items():
                    paragraph += f"{column_name}: {value}\n"
                paragraphs.append(paragraph.strip())
            return paragraphs

class XLSXFileLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def getParagraphs(self):
        wb = openpyxl.load_workbook(self.file_path)
        ws = wb.active
        paragraphs = []
        header = [cell.value for cell in ws[1]]
        for row in ws.iter_rows(min_row=2, values_only=True):
            paragraph = ""
            for column_name, value in zip(header, row):
                paragraph += f"{column_name}: {value}\n"
            paragraphs.append(paragraph.strip())
        return paragraphs

class VectorDBManager:
    def __init__(self, embedding_func):
        self.embedding_func = embedding_func

    def create_vector_db_from_folder(self, folder_path, db_path):
        """通过调用 AI 模型，创建向量数据库"""
        folder_name = os.path.basename(folder_path.rstrip(os.sep))
        print(f"Creating vector database for folder: {folder_name}")

        # 初始化向量数据库
        vector_db = MyVectorDBConnector(folder_name, self.embedding_func, db_path)

        loaders = {
            ".pdf": PDFFileLoader,
            ".doc": DocFileLoader,
            ".docx": DocFileLoader,
            ".ppt": PPTFileLoader,
            ".pptx": PPTFileLoader,
            '.csv': CSVFileLoader,
            '.xlsx': XLSXFileLoader,
        }

        for root, _, files in os.walk(folder_path):
            for file_name in files:
                file_path = os.path.join(root, file_name)
                file_ext = os.path.splitext(file_name)[1].lower()

                if file_ext in loaders:
                    try:
                        # 加载文件内容
                        paragraphs = loaders[file_ext](file_path).getParagraphs()
                        file_text = "\n".join(paragraphs)

                        # 检查是否为表格文件（csv或xlsx）
                        if file_ext in ['.csv', '.xlsx']:
                            # 按表格行分割
                            chunks = [para for para in paragraphs if para.strip()]
                        else:
                            # 其他文件类型按原逻辑处理
                            file_text = "\n".join(paragraphs)
                            chunks = chunk_text(file_text, chunk_size=1000, overlap=200)

                        # 遍历切片并处理
                        for i, chunk in enumerate(chunks):
                            try:
                                result = process_chunk_with_ai(chunk, f"{file_path}_chunk_{i}", self.embedding_func, vector_db)
                                print(result)
                                print(f"Successfully processed chunk {i + 1}/{len(chunks)} for {file_name}")
                            except Exception as e:
                                print(f"Error processing chunk {i + 1}/{len(chunks)} for {file_name}: {e}")

                    except Exception as e:
                        print(f"Error processing {file_name}: {e}")
                else:
                    print(f"Unsupported file type: {file_name}")

        return vector_db
     

db_path = '/home/student/zlh/GuardBot-main/chroma_data'

directory_to_process = "/home/student/zlh/GuardBot-main/new_doc_db"  # 替换为你的目录路径
dirs = []
for file in os.scandir(directory_to_process):
    # 检查是否为一级目录
    if file.is_dir():
        print(f"Processing directory: {file.name}")
        
        # 创建 VectorDBManager 实例
        db_manager = VectorDBManager(get_embeddings)
        
        # 调用方法对当前目录创建向量数据库
        try:
            vector_db = db_manager.create_vector_db_from_folder(file.path, db_path)
            print(f"Vector database created for directory: {file.name}")
        except Exception as e:
            print(f"Error processing directory {file.name}: {e}")

# db_path = "/home/student/zlh/GuardBot-main/chroma_data"  # 替换为你的数据库路径
# # 合法化后的集合名称
# valid_collection_name = chinese_to_valid_name('保险法')  # 替换为你的集合名称

# # 初始化客户端
# chroma_client = PersistentClient(path=db_path)

# # 删除集合
# chroma_client.delete_collection(name=valid_collection_name)


